"""
QQ Bot 工具系统
提供 bot 可以主动调用的工具：
  web_search / add_memory / search_memory / update_soul
  read_file / write_file / list_files（workspace 文件操作，write 需要 ADMIN+）
"""

import asyncio
import json
import logging
import re
import sqlite3
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Tuple

if TYPE_CHECKING:
    from .permissions import PermLevel as _PermLevel

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────
# 工具 Schema（OpenAI function-calling 格式）
# ──────────────────────────────────────────────────────────────

TOOL_SCHEMAS = [
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "上网搜索信息。当需要查找最新新闻、实时数据、不确定的事实时使用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "搜索关键词，尽量简洁准确"
                    }
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "add_memory",
            "description": (
                "记住某件重要的事情。"
                "scope=\"person\"（默认）：记录关于**这个人**的信息（姓名/喜好/职业/重要事件），下次见到他/她时能想起来。"
                "scope=\"global\"：记录关于**小萌自身**的规则、设定、主人的要求等，所有对话都生效。如：称呼规则、行为准则、主人指定的偏好。"
                "scope=\"knowledge\"：记录关于**外部世界**的知识，如动漫设定、乐队成员、历史事件、技术概念等。"
                "scope=\"conversation\"：记录关于**当前对话主题**的信息，仅在本聊天窗口有效。"
                "【重要】对方一旦透露了个人信息（哪怕只是名字、年级、喜好），立即用 scope=\"person\" 记下来！"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": "【必填，字段名必须是 content】要记住的内容，用一句话描述清楚。必须包含：谁说/在哪看到 + 关于谁 + 说了什么。如「A说B喜欢编程」「B自称叫B」「网上查到 Ave Mujica 成员有…」"
                    },
                    "about": {
                        "type": "string",
                        "description": "scope=\"person\" 时指定这条记忆是关于谁的。取值为 system prompt 已知人物列表中列出的 identity。不填则默认为当前对话对象。"
                    },
                    "scope": {
                        "type": "string",
                        "description": "person=关于这个人；global=关于小萌自身的规则/主人的要求；knowledge=关于外部世界的知识；conversation=关于当前对话",
                        "enum": ["person", "global", "knowledge", "conversation"]
                    },
                    "importance": {
                        "type": "integer",
                        "description": "重要程度 1-3，1=普通，2=重要，3=非常重要",
                        "enum": [1, 2, 3]
                    }
                },
                "required": ["content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_memory",
            "description": "搜索自己的记忆，回忆之前聊过的事情或这个人说过的话。会同时搜索关于这个人的记忆和当前对话的记忆。",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "想要搜索的关键词或描述"
                    }
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "update_soul",
            "description": (
                "更新自己的 SOUL.md 进化区，记录真实的成长和新发现。"
                "适合用在：遇到了以前没想过的问题、某次对话让你有了新的认知、发现自己喜欢/不喜欢某件事。"
                "不要强行写，要真实。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": "以「- 」开头的一两句话，说清楚是什么、为什么值得写进来"
                    }
                },
                "required": ["content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "recall_conversations",
            "description": (
                "回忆近期和其他人聊过的内容（跨群跨私聊）。仅主人和管理员可用。"
                "主人问「最近和谁聊了什么」「群里有没有人提过xxx」时使用。"
                "可以按关键词搜索，也可以查最近的所有对话记录。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "要搜索的关键词（为空时返回最近的对话记录）"
                    },
                    "limit": {
                        "type": "integer",
                        "description": "返回最多多少条消息，默认 20",
                    }
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": (
                "读取 workspace（data/）下的任意文件。"
                "可以读 persona/SOUL.md、memory/owner.md、identity_links.json、skills/ 下的技能文件等。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "相对于 data/ 的路径，如 'persona/SOUL.md' 或 'identity_links.json'"
                    }
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": (
                "写入或修改 workspace（data/）下的文件。需要管理员或主人权限。"
                "可以用来：更新 identity_links.json 把某个 QQ 号关联到一个人、"
                "修改某人的记忆文件（memory/xxx.md）、更新 persona 文件等。"
                "修改 identity_links.json 后身份映射立即生效。"
                "path 留空时自动在 notes/ 下创建带时间戳的新文件。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "相对于 data/ 的路径；留空则自动生成文件名（notes/时间戳.md）"
                    },
                    "content": {
                        "type": "string",
                        "description": "要写入的完整内容（overwrite 模式）或追加的内容（append 模式）"
                    },
                    "mode": {
                        "type": "string",
                        "description": "overwrite（覆盖，默认）或 append（在末尾追加）",
                        "enum": ["overwrite", "append"]
                    }
                },
                "required": ["path", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "列出 workspace（data/）某目录下的文件列表。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "相对于 data/ 的目录路径，空字符串表示根目录"
                    }
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": (
                "删除 workspace（data/）下的文件或空目录。需要管理员或主人权限。"
                "可以用来：删除过期的 skill、清理不需要的笔记、移除错误的记忆文件等。"
                "只能删 data/ 内的文件，不能删系统文件。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "相对于 data/ 的路径，如 'skills/old_skill.md' 或 'notes/draft.md'"
                    }
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_voice",
            "description": (
                "合成并发送语音消息。适合：用户明确要求语音回复、朗读诗词/歌词、"
                "简短情感化的表达需要语音时。调用后会直接发送语音，无需再重复文字内容。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {
                        "type": "string",
                        "description": "要合成为语音的文字内容，建议100字以内"
                    }
                },
                "required": ["text"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": (
                "在服务器上执行 shell 命令。**仅主人可用。**"
                "可以用来：管理服务，查看日志，安装依赖，重启程序，查看系统状态等。"
                "返回命令的 stdout 和 stderr。超时 60 秒自动终止。"
                "工作目录为 XiaoMeng 项目根目录。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {
                        "type": "string",
                        "description": "要执行的 shell 命令，如 'systemctl status nginx' 或 'pip install requests'"
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "超时秒数，默认 30，最大 60"
                    }
                },
                "required": ["command"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_message",
            "description": (
                "在任务执行过程中向用户发送一条消息，用于告知进度、请求确认或询问信息。"
                "如果需要等待用户回复，设置 wait_reply=true，此时任务会暂停，"
                "收到用户回复后自动继续执行。"
                "【重要】执行长任务时，先用 send_message 告诉用户「收到，正在处理」，"
                "然后继续执行工具，最后再用 send_message 发送最终结果。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "message": {
                        "type": "string",
                        "description": "要发送给用户的消息内容"
                    },
                    "wait_reply": {
                        "type": "boolean",
                        "description": "是否等待用户回复后再继续。需要用户确认或提供信息时设为 true"
                    }
                },
                "required": ["message"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_file",
            "description": (
                "将 workspace（data/）下的文件发送给当前聊天对象（群或私聊）。"
                "适合：发送生成的文档、题目文件、笔记、图片等。"
                "可附带一条说明消息一起发送。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "相对于 data/ 的文件路径，如 'skills/daily-quiz/questions.md'"
                    },
                    "caption": {
                        "type": "string",
                        "description": "发送文件前附带的说明消息（可选）"
                    }
                },
                "required": ["path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "reload_skills",
            "description": (
                "重新加载 skills 目录下的所有技能文件。"
                "当你用 write_file 写入新 skill（skills/ 目录下的 .md 文件）后，"
                "必须调用此工具使其立即生效。"
                "调用后所有新写入或修改的技能文件会被重新扫描加载。"
            ),
            "parameters": {
                "type": "object",
                "properties": {},
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "send_to",
            "description": (
                "主动向指定群或用户发送一条消息（仅限内心世界/主人权限使用）。"
                "group_id 和 user_id 二选一，不能同时填。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "message":  {"type": "string",  "description": "要发送的文字内容"},
                    "group_id": {"type": "integer", "description": "目标群号（与 user_id 二选一）"},
                    "user_id":  {"type": "integer", "description": "目标用户 QQ 号（与 group_id 二选一）"},
                },
                "required": ["message"],
            },
        },
        "_meta": {"min_user_level": 2, "risk": "SENSITIVE"},
    },
    {
        "type": "function",
        "function": {
            "name": "add_reminder",
            "description": (
                "设置定时提醒。支持自然语言时间表达，如：\n"
                "- 一次性：'10分钟后'、'今晚8点'、'明天下午5点'、'后天早上9点'\n"
                "- 每日重复：'每天早上7点'、'每日晚上10点'\n"
                "到时间后会私信或在群内提醒。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "time_str": {
                        "type": "string",
                        "description": "自然语言时间，如'明天下午5点'、'10分钟后'、'每天早上8点'"
                    },
                    "message": {
                        "type": "string",
                        "description": "提醒内容"
                    },
                    "remind_in_group": {
                        "type": "boolean",
                        "description": "是否在群里提醒（默认 false，发私信）"
                    }
                },
                "required": ["time_str", "message"]
            }
        }
    },
]

# 调用工具时展示给用户看的进度消息（None = 不展示）
TOOL_PROGRESS_MSG = {
    "web_search":           lambda args: f"让小萌查一下「{args.get('query', '')}」~",
    "send_voice":           lambda args: "（小萌合成语音中，稍等一下~）",
    "run_command":          lambda args: f"（小萌执行命令中：`{args.get('command', '')}` ~）",
    "write_file":           lambda args: f"（小萌写入 {args.get('path') or args.get('file_path', '文件')} ~）",
    "delete_file":          lambda args: f"（小萌删除 {args.get('path') or args.get('file_path', '文件')} ~）",
    "add_memory":           lambda args: "让小萌记下来",
    "search_memory":        lambda args: None,
    "update_soul":          lambda args: None,
    "recall_conversations": lambda args: None,
    "read_file":            lambda args: None,
    "list_files":           lambda args: None,
    "reload_skills":        lambda args: "（小萌重新加载技能中~）",
    "send_file":            lambda args: f"（小萌发送文件 {args.get('path', '')} ~）",
    "add_reminder":         lambda args: f"（小萌设置提醒：{args.get('message', '')}~）",
}


# ──────────────────────────────────────────────────────────────
# 工具执行器
# ──────────────────────────────────────────────────────────────

class QQToolExecutor:
    """执行工具调用，返回结果字符串。"""

    def __init__(
        self,
        db_path: str,
        soul_path: Path,
        data_dir: Path,     # workspace 根目录（data/），沙箱范围
        session_key: str,
        sender_qq: int = 0,
        level=None,         # PermLevel，write_file 权限检查用
        identity: str = "", # canonical 身份名，记忆文件用
        proxy: str = "",    # web_search 代理
        napcat=None,        # NapCatClient，send_voice 需要
        tts=None,           # SoVITSTTS，send_voice 需要
        target_id: int = 0, # 发送目标（group_id 或 user_id）
        is_private: bool = False,
        task=None,          # AsyncTask，send_message + wait_reply 需要
        plugin_manager=None,  # PluginManager，插件工具分发
    ):
        self._db_path = db_path
        self._soul_path = soul_path
        self._data_dir = data_dir.resolve()
        self._memory_dir = data_dir / "memory"
        self._memory_md_path = soul_path.parent / "MEMORY.md"
        self._session_key = session_key
        self._sender_qq = sender_qq
        self._level = level
        self._identity = identity or (f"user_{sender_qq}" if sender_qq else "unknown")
        self._user_key = f"user:{sender_qq}" if sender_qq else session_key
        self._proxy = proxy
        self._napcat = napcat
        self._tts = tts
        self._target_id = target_id
        self._is_private = is_private
        self._task = task
        self._plugin_manager = plugin_manager

    async def execute(self, tool_name: str, arguments: Dict[str, Any]) -> str:
        if tool_name == "web_search":
            return await self._web_search(arguments.get("query", ""))
        elif tool_name == "add_memory":
            # LLM 有时用各种字段名代替 content，逐一尝试
            content = (arguments.get("content")
                       or arguments.get("summary")
                       or arguments.get("text")
                       or arguments.get("note")
                       or arguments.get("value")
                       or arguments.get("fact")
                       or arguments.get("info", ""))
            # key+value 组合：拼成 "key: value"
            if not content and arguments.get("key") and arguments.get("value"):
                content = f"{arguments['key']}: {arguments['value']}"
            scope = arguments.get("scope", "person")
            if scope not in ("person", "conversation", "knowledge", "global"):
                scope = "person"
            about = arguments.get("about") or ""
            return self._add_memory(content, int(arguments.get("importance", 1)), scope, about)
        elif tool_name == "search_memory":
            return self._search_memory(arguments.get("query", ""))
        elif tool_name == "update_soul":
            content = arguments.get("content", "")
            # LLM 有时传 updates（列表）而不是 content（字符串），兼容处理
            if not content:
                updates = arguments.get("updates", [])
                if isinstance(updates, list):
                    parts = []
                    for u in updates:
                        if isinstance(u, dict):
                            parts.append(u.get("affirmation") or u.get("content") or u.get("text") or str(u))
                        elif isinstance(u, str):
                            parts.append(u)
                    content = "\n".join(p for p in parts if p)
                elif isinstance(updates, str):
                    content = updates
            return self._update_soul(content)
        elif tool_name == "recall_conversations":
            return self._recall_conversations(
                arguments.get("query", ""),
                int(arguments.get("limit", 20)),
            )
        elif tool_name == "read_file":
            path = arguments.get("path") or arguments.get("file_path", "")
            return self._read_file(path)
        elif tool_name == "write_file":
            # LLM 有时会传 file_path 而不是 path，兼容两种写法
            path = arguments.get("path") or arguments.get("file_path", "")
            return self._write_file(
                path,
                arguments.get("content", ""),
                arguments.get("mode", "overwrite"),
            )
        elif tool_name == "list_files":
            return self._list_files(arguments.get("path", ""))
        elif tool_name == "delete_file":
            path = arguments.get("path") or arguments.get("file_path", "")
            return self._delete_file(path)
        elif tool_name == "send_voice":
            return await self._send_voice(arguments.get("text", ""))
        elif tool_name == "send_message":
            return await self._send_message(
                arguments.get("message", ""),
                arguments.get("wait_reply", False),
            )
        elif tool_name == "run_command":
            return await self._run_command(
                arguments.get("command", ""),
                int(arguments.get("timeout", 30)),
            )
        elif tool_name == "send_file":
            path = arguments.get("path") or arguments.get("file_path", "")
            return await self._send_file(path, arguments.get("caption", ""))
        elif tool_name == "send_to":
            msg = arguments.get("message", "")
            gid = arguments.get("group_id")
            uid = arguments.get("user_id")
            if not msg:
                return "参数错误：message 不能为空"
            if gid:
                await self._napcat.send_group_msg(int(gid), msg)
                return f"已发送到群 {gid}"
            elif uid:
                await self._napcat.send_private_msg(int(uid), msg)
                return f"已发送到用户 {uid}"
            else:
                return "参数错误：group_id 和 user_id 至少填一个"
        elif tool_name == "reload_skills":
            return self._reload_skills()
        elif tool_name == "add_reminder":
            return self._add_reminder(
                arguments.get("time_str", ""),
                arguments.get("message", ""),
                bool(arguments.get("remind_in_group", False)),
            )
        elif self._plugin_manager is not None:
            ctx = {
                "session_key": self._session_key,
                "sender_qq": self._sender_qq,
                "level": self._level,
                "identity": self._identity,
                "napcat": self._napcat,
                "data_dir": self._data_dir,
                "is_private": self._is_private,
                "target_id": self._target_id,
                "task": self._task,
            }
            return await self._plugin_manager.dispatch_tool_call(tool_name, arguments, ctx)
        else:
            return f"未知工具: {tool_name}"

    # ── web_search ──────────────────────────────────

    async def _web_search(self, query: str) -> str:
        if not query.strip():
            return "搜索词不能为空"
        try:
            from ddgs import DDGS
            import asyncio
            results = await asyncio.to_thread(_ddg_search, query, 4, self._proxy)
            if not results:
                return "没找到相关结果"
            lines = []
            for r in results:
                title = r.get("title", "")
                body = r.get("body", "")[:150]
                href = r.get("href", "")
                lines.append(f"**{title}**\n{body}\n{href}")
            return "\n\n".join(lines)
        except Exception as e:
            logger.warning(f"web_search 失败: {e}")
            return f"搜索失败: {e}"

    # ── add_memory ───────────────────────────────────

    def _add_memory(self, content: str, importance: int = 1, scope: str = "person", about: str = "") -> str:
        if not content.strip():
            return "内容不能为空"
        content = content.strip()
        now = datetime.now()
        ts = now.strftime("%Y-%m-%d")

        if scope == "global":
            try:
                if not self._memory_md_path.exists():
                    self._memory_md_path.write_text("# 小萌的重要记忆\n\n", encoding="utf-8")
                with open(self._memory_md_path, "a", encoding="utf-8") as f:
                    f.write(f"- [{ts}] {content}\n")
            except Exception as e:
                return f"全局记忆写入失败: {e}"
            try:
                conn = sqlite3.connect(self._db_path)
                conn.execute(
                    "INSERT INTO long_term_memory (session_key, content, memory_type, importance, created_at) VALUES (?, ?, 'global', ?, ?)",
                    ("global", content, importance, now.isoformat()),
                )
                conn.commit()
                conn.close()
            except Exception:
                pass
            return f"记住了（全局）：{content}"

        if scope == "knowledge":
            # 通用知识：写 data/memory/knowledge.md + SQLite key=knowledge
            try:
                knowledge_file = self._memory_dir / "knowledge.md"
                if not knowledge_file.exists():
                    knowledge_file.write_text("# 小萌学到的知识\n\n", encoding="utf-8")
                with open(knowledge_file, "a", encoding="utf-8") as f:
                    f.write(f"- [{ts}] {content}\n")
            except Exception as e:
                return f"知识写入失败: {e}"
            try:
                conn = sqlite3.connect(self._db_path)
                conn.execute(
                    "INSERT INTO long_term_memory (session_key, content, memory_type, importance, created_at) VALUES (?, ?, 'knowledge', ?, ?)",
                    ("knowledge", content, importance, now.isoformat()),
                )
                conn.commit()
                conn.close()
            except Exception:
                pass
            return f"记住了（知识库）：{content}"

        elif scope == "person" and self._sender_qq:
            target_identity = about or self._identity
            try:
                mem_file = self._memory_dir / f"{target_identity}.md"
                line = f"- [{ts}] {content}\n"
                with open(mem_file, "a", encoding="utf-8") as f:
                    f.write(line)
            except Exception as e:
                return f"记忆写入文件失败: {e}"
            try:
                user_key = f"user:{self._sender_qq}"
                conn = sqlite3.connect(self._db_path)
                conn.execute(
                    "INSERT INTO long_term_memory (session_key, content, memory_type, importance, created_at) VALUES (?, ?, 'person', ?, ?)",
                    (user_key, content, importance, now.isoformat()),
                )
                conn.commit()
                conn.close()
            except Exception:
                pass
        else:
            # 写到当前会话（群/私聊）的 SQLite 记忆
            try:
                conn = sqlite3.connect(self._db_path)
                conn.execute(
                    "INSERT INTO long_term_memory (session_key, content, memory_type, importance, created_at) VALUES (?, ?, 'explicit', ?, ?)",
                    (self._session_key, content, importance, now.isoformat()),
                )
                conn.commit()
                conn.close()
            except Exception as e:
                return f"记忆失败: {e}"

        return f"记住了：{content}"

    # ── search_memory ────────────────────────────────

    def _search_memory(self, query: str) -> str:
        if not query.strip():
            return "搜索词不能为空"
        like = f"%{query.strip()}%"
        lines = []
        try:
            conn = sqlite3.connect(self._db_path)
            conn.row_factory = sqlite3.Row
            c = conn.cursor()

            # 1. 长期记忆：当前会话 + 用户个人记忆 + 知识库 + 全局设定
            keys = list({self._session_key, self._user_key, "knowledge", "global"})
            placeholders = ",".join("?" * len(keys))
            c.execute(
                f"""SELECT content, session_key, memory_type, importance, created_at
                   FROM long_term_memory
                   WHERE session_key IN ({placeholders}) AND content LIKE ?
                   ORDER BY importance DESC, id DESC LIMIT 8""",
                (*keys, like),
            )
            for r in [dict(r) for r in c.fetchall()]:
                ts = (r.get("created_at") or "")[:10]
                if r["session_key"] == "knowledge":
                    tag = "（知识库）"
                elif r["session_key"] == "global":
                    tag = "（全局设定）"
                elif r["session_key"] == self._user_key and self._user_key != self._session_key:
                    tag = "（关于ta）"
                else:
                    tag = ""
                lines.append(f"[{ts}]{tag} {r['content']}")

            # 2. 原始消息表：搜当前会话里小萌自己说过的话
            c.execute(
                """SELECT content, created_at FROM messages
                   WHERE session_key=? AND role='assistant' AND content LIKE ?
                   ORDER BY id DESC LIMIT 5""",
                (self._session_key, like),
            )
            for r in [dict(r) for r in c.fetchall()]:
                ts = (r.get("created_at") or "")[:10]
                lines.append(f"[{ts}]（小萌说过） {r['content'][:120]}")

            conn.close()
        except Exception as e:
            return f"记忆搜索失败: {e}"

        if not lines:
            return "没找到相关记忆"
        return "\n".join(lines)

    # ── update_soul ──────────────────────────────────

    def _update_soul(self, content: str) -> str:
        if not content.strip():
            return "内容不能为空"
        try:
            soul_text = self._soul_path.read_text(encoding="utf-8")
            evolving_end = soul_text.find("<!-- EVOLVING_END")
            if evolving_end == -1:
                return "SOUL.md 格式不对，找不到进化区标记"

            content = content.strip()
            if not content.startswith("- "):
                content = "- " + content

            now = datetime.now().strftime("%Y-%m-%d")
            insert = f"\n{content}  <!-- {now} -->\n"
            new_soul = soul_text[:evolving_end] + insert + soul_text[evolving_end:]
            self._soul_path.write_text(new_soul, encoding="utf-8")
            return f"SOUL.md 进化区已更新：{content}"
        except Exception as e:
            return f"更新失败: {e}"

    # ── recall_conversations ──────────────────────────

    def _recall_conversations(self, query="", limit: int = 20) -> str:
        query = str(query) if query else ""  # 防止 LLM 传 int
        # PermLevel.ADMIN == 1，用整数比较避免循环导入
        if self._level is None or int(self._level) < 1:
            return "没有权限查看对话记录（需要管理员或主人）"
        limit = max(1, min(limit, 100))
        try:
            conn = sqlite3.connect(self._db_path)
            conn.row_factory = sqlite3.Row
            c = conn.cursor()
            if query.strip():
                # 拆词：识别群/私聊类型词，剩余做内容关键词搜索
                raw_words = query.strip().split()
                session_prefix = None   # "group:" 或 "private:"
                kw_words = []
                for w in raw_words:
                    if w in ("群", "群聊", "group"):
                        session_prefix = "group:%"
                    elif w in ("私", "私聊", "私信", "private"):
                        session_prefix = "private:%"
                    else:
                        kw_words.append(w)

                # 构建 WHERE 子句
                where_parts = []
                params: list = []

                # session_key 过滤（可选）
                if session_prefix:
                    where_parts.append("session_key LIKE ?")
                    params.append(session_prefix)

                # 内容关键词（按词 OR，每词 content+sender_name）
                if kw_words:
                    kw_conds = []
                    for w in kw_words:
                        like = f"%{w}%"
                        kw_conds.append("(content LIKE ? OR sender_name LIKE ?)")
                        params.extend([like, like])
                    where_parts.append("(" + " OR ".join(kw_conds) + ")")

                if not where_parts:
                    where = "1"   # 不会走到这里，但做兜底
                elif len(where_parts) == 1:
                    where = where_parts[0]
                else:
                    where = " AND ".join(where_parts)

                params.append(limit)
                c.execute(
                    f"""SELECT session_key, role, sender_name, content, created_at
                       FROM messages
                       WHERE {where}
                       ORDER BY id DESC LIMIT ?""",
                    params,
                )
            else:
                c.execute(
                    """SELECT session_key, role, sender_name, content, created_at
                       FROM messages
                       ORDER BY id DESC LIMIT ?""",
                    (limit,),
                )
            rows = [dict(r) for r in c.fetchall()]
            conn.close()
            if not rows:
                return "没有找到相关对话记录"
            # 倒序排列（最旧在前），按 session 分组显示
            rows = list(reversed(rows))
            lines = []
            last_session = None
            for r in rows:
                sk = r["session_key"]
                if sk != last_session:
                    # 把 session_key 翻译成更可读的标题
                    if sk.startswith("group:"):
                        header = f"群聊 {sk[6:]}"
                    elif sk.startswith("private:"):
                        header = f"私聊 {sk[8:]}"
                    else:
                        header = sk
                    lines.append(f"\n# {header}：")
                    last_session = sk
                ts = (r.get("created_at") or "")[:16]
                name = r.get("sender_name") or ("小萌" if r["role"] == "assistant" else "?")
                lines.append(f"  [{ts}] {name}: {r['content'][:80]}")
            return "\n".join(lines)
        except Exception as e:
            return f"查询失败: {e}"

    # ── read_file ─────────────────────────────────────

    def _read_file(self, path: str) -> str:
        if not path:
            return "路径不能为空"
        try:
            target = self._resolve_workspace_path(path)
            if not target.exists():
                return f"文件不存在: {path}"
            if target.is_dir():
                return f"{path} 是目录，请用 list_files"
            content = self._extract_text(target)
            if len(content) > 4000:
                content = content[:4000] + f"\n…（文件还有 {len(content)-4000} 字符，已截断）"
            return content
        except ValueError as e:
            return f"路径错误: {e}"
        except Exception as e:
            return f"读取失败: {e}"

    @staticmethod
    def _extract_text(path: Path) -> str:
        """从各种文件格式提取纯文本。"""
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            import fitz  # pymupdf
            doc = fitz.open(str(path))
            pages = [doc[i].get_text() for i in range(len(doc))]
            doc.close()
            return "\n".join(pages).strip()

        if suffix in (".docx",):
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        if suffix in (".pptx",):
            from pptx import Presentation
            prs = Presentation(str(path))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
            return "\n".join(lines).strip()

        if suffix in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        lines.append("\t".join(cells))
            wb.close()
            return "\n".join(lines).strip()

        # 其他格式当纯文本读
        return path.read_text(encoding="utf-8", errors="replace")

    # ── write_file ────────────────────────────────────

    def _write_file(self, path: str, content: str, mode: str = "overwrite") -> str:
        # 权限检查：只有管理员和主人可以写文件（PermLevel.ADMIN == 1）
        if self._level is None or int(self._level) < 1:
            return "没有权限修改文件（需要管理员或主人）"
        # 路径为空 → 自动生成文件名（notes/YYYYMMDD_HHMMSS.md）
        if not path:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            path = f"notes/{ts}.md"
        try:
            target = self._resolve_workspace_path(path)
            target.parent.mkdir(parents=True, exist_ok=True)
            if mode == "append":
                with open(target, "a", encoding="utf-8") as f:
                    f.write(content)
            else:
                target.write_text(content, encoding="utf-8")
            logger.info(f"write_file: {path} ({mode}) by {self._sender_qq}")
            # 如果写入了 skills/ 目录，自动重新加载技能注册中心
            if path.strip("/").startswith("skills"):
                self._reload_skills()
                return f"已写入 {path}（技能已自动加载生效）"
            return f"已写入 {path}"
        except ValueError as e:
            return f"路径错误: {e}"
        except Exception as e:
            return f"写入失败: {e}"

    # ── list_files ────────────────────────────────────

    def _list_files(self, path: str = "") -> str:
        try:
            target = self._resolve_workspace_path(path) if path else self._data_dir
            if not target.exists():
                return f"目录不存在: {path or '(root)'}"
            if not target.is_dir():
                return f"{path} 不是目录"
            items = []
            for item in sorted(target.iterdir()):
                if item.is_dir():
                    items.append(f"[dir]  {item.name}/")
                else:
                    size = item.stat().st_size
                    items.append(f"[file] {item.name}  ({size} bytes)")
            return "\n".join(items) if items else "（空目录）"
        except ValueError as e:
            return f"路径错误: {e}"
        except Exception as e:
            return f"列出失败: {e}"

    # ── delete_file ───────────────────────────────────

    def _delete_file(self, path: str) -> str:
        if self._level is None or int(self._level) < 1:
            return "没有权限删除文件（需要管理员或主人）"
        if not path:
            return "路径不能为空"
        # 保护核心文件不被误删
        PROTECTED = {"persona/SOUL.md", "identity_links.json", "qq_bot.db"}
        if path.strip("/") in PROTECTED:
            return f"禁止删除核心文件: {path}"
        try:
            target = self._resolve_workspace_path(path)
            if not target.exists():
                return f"文件不存在: {path}"
            if target.is_dir():
                # 只允许删空目录
                if any(target.iterdir()):
                    return f"目录非空，无法删除: {path}（请先删除里面的文件）"
                target.rmdir()
            else:
                target.unlink()
            logger.info(f"delete_file: {path} by {self._sender_qq}")
            if path.strip("/").startswith("skills"):
                self._reload_skills()
                return f"已删除 {path}（技能注册已同步更新）"
            return f"已删除 {path}"
        except ValueError as e:
            return f"路径错误: {e}"
        except Exception as e:
            return f"删除失败: {e}"

    # ── send_voice ────────────────────────────────────

    async def _send_voice(self, text: str) -> str:
        import base64
        if not text.strip():
            return "语音内容不能为空"
        if not self._tts:
            return "TTS 服务未配置，无法发送语音"
        if not self._napcat or not self._target_id:
            return "发送目标未配置，无法发送语音"
        audio_bytes = await self._tts.synthesize(text)
        if not audio_bytes:
            return "语音合成失败，请稍后再试"
        b64 = base64.b64encode(audio_bytes).decode()
        cq = f"[CQ:record,file=base64://{b64}]"
        try:
            if self._is_private:
                await self._napcat.send_private_msg(self._target_id, cq)
            else:
                await self._napcat.send_group_msg(self._target_id, cq)
            return "✅ 语音消息已发送"
        except Exception as e:
            return f"语音发送失败: {e}"

    # ── run_command ───────────────────────────────────

    async def _run_command(self, command: str, timeout: int = 30) -> str:
        if self._level is None or int(self._level) < 2:  # 2 == PermLevel.OWNER
            return "没有权限执行命令（仅主人可用）"
        if not command.strip():
            return "命令不能为空"
        timeout = max(1, min(timeout, 60))
        import asyncio, subprocess
        # 工作目录：项目根（data/ 的上一级）
        cwd = str(self._data_dir.parent)
        logger.info(f"run_command by {self._sender_qq}: {command!r}")
        try:
            proc = await asyncio.create_subprocess_shell(
                command,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                cwd=cwd,
            )
            try:
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
            except asyncio.TimeoutError:
                proc.kill()
                return f"命令超时（>{timeout}s）已终止"
            out = stdout.decode("utf-8", errors="replace").strip()
            err = stderr.decode("utf-8", errors="replace").strip()
            rc = proc.returncode
            parts = []
            if out:
                parts.append(out)
            if err:
                parts.append(f"[stderr]\n{err}")
            result = "\n".join(parts) if parts else "（无输出）"
            # 回复太长时截断
            if len(result) > 3000:
                result = result[:3000] + f"\n…（输出过长，已截断，共 {len(result)} 字符）"
            return f"返回码: {rc}\n{result}"
        except Exception as e:
            return f"执行失败: {e}"

    # ── workspace 路径解析（防穿越）─────────────────────

    async def _send_message(self, message: str, wait_reply: bool = False) -> str:
        if not message.strip():
            return "消息内容不能为空"
        if not self._napcat:
            return "send_message: napcat 不可用"
        try:
            if self._is_private:
                await self._napcat.send_private_msg(self._target_id, message)
            else:
                await self._napcat.send_group_msg(self._target_id, message)
        except Exception as e:
            logger.error(f"send_message 发送失败: {e}")
            return f"消息发送失败: {e}"

        if wait_reply and self._task is not None:
            user_reply = await self._task.wait_for_user(f"已发送: {message[:50]}...")
            if user_reply is None:
                return "等待用户回复超时"
            return f"用户回复: {user_reply}"
        else:
            return f"已发送消息: {message[:100]}"

    # ── send_file ─────────────────────────────────────

    async def _send_file(self, path: str, caption: str = "") -> str:
        import base64 as _b64
        if not path:
            return "路径不能为空"
        if not self._napcat or not self._target_id:
            return "发送目标未配置，无法发送文件"
        try:
            target = self._resolve_workspace_path(path)
            if not target.exists():
                return f"文件不存在: {path}"
            if target.is_dir():
                return f"{path} 是目录，请指定具体文件"
            if caption:
                try:
                    if self._is_private:
                        await self._napcat.send_private_msg(self._target_id, caption)
                    else:
                        await self._napcat.send_group_msg(self._target_id, caption)
                except Exception:
                    pass
            # NapCat 运行在容器内，无法访问宿主机 file:// 路径
            # 用 base64:// 直接传文件内容，绕过路径可访问性问题
            file_b64 = _b64.b64encode(target.read_bytes()).decode()
            file_ref = f"base64://{file_b64}"
            file_name = target.name
            if self._is_private:
                resp = await self._napcat.call_api("upload_private_file", {
                    "user_id": self._target_id,
                    "file": file_ref,
                    "name": file_name,
                }, timeout=120.0)
            else:
                resp = await self._napcat.call_api("upload_group_file", {
                    "group_id": self._target_id,
                    "file": file_ref,
                    "name": file_name,
                }, timeout=120.0)
            if resp.get("retcode", 0) != 0:
                return f"发送失败 (retcode={resp['retcode']}): {resp.get('msg', '')}"
            return f"✅ 文件已发送: {file_name}"
        except ValueError as e:
            return f"路径错误: {e}"
        except Exception as e:
            return f"发送失败: {e}"

    def _resolve_workspace_path(self, path: str) -> Path:
        resolved = (self._data_dir / path).resolve()
        if not str(resolved).startswith(str(self._data_dir)):
            raise ValueError(f"路径越界，只能访问 data/ 目录内的文件: {path}")
        return resolved

    def _reload_skills(self) -> str:
        """重新加载 skills/ 下所有技能文件，并重新注册各插件的 SKILL.md。"""
        try:
            from .skills import SkillRegistry, SkillLoader
            registry = SkillRegistry.get_instance()
            count = registry.reload()

            # 同步重新注册插件 SKILL.md
            if self._plugin_manager is not None:
                for _name, _plugin in self._plugin_manager.initialized.items():
                    _skill_md = _plugin.plugin_dir / "SKILL.md"
                    if _skill_md.exists():
                        _loader = SkillLoader(str(_plugin.plugin_dir))
                        _skill_def = _loader.load_single("SKILL.md")
                        if _skill_def:
                            registry.register(_skill_def, override=True)
                            count += 1

            if count == 0:
                return "没有找到任何技能文件"
            names = ", ".join(s.name for s in registry.list_all())
            return f"已重新加载 {count} 个技能：{names}"
        except Exception as e:
            logger.error(f"reload_skills 失败: {e}")
            return f"重新加载技能失败: {e}"

    @staticmethod
    def _reload_skills_registry() -> str:
        """向后兼容保留，内部调 SkillRegistry.reload()。"""
        try:
            from .skills import SkillRegistry
            registry = SkillRegistry.get_instance()
            count = registry.reload()
            names = ", ".join(s.name for s in registry.list_all())
            return f"已重新加载 {count} 个技能：{names}"
        except Exception as e:
            return f"重新加载技能失败: {e}"

    # ── add_reminder ─────────────────────────────────

    def _add_reminder(self, time_str: str, message: str, remind_in_group: bool = False) -> str:
        parsed = self._parse_reminder_time(time_str)
        if parsed is None:
            return f"小萌没看懂这个时间表达：「{time_str}」，可以试试「明天下午5点」或「10分钟后」~"
        task_type, run_at = parsed
        import uuid, json as _json
        task_id = f"reminder_{uuid.uuid4().hex[:8]}"
        if remind_in_group and self._target_id and not self._is_private:
            target = {"group_id": self._target_id}
        elif self._sender_qq:
            target = {"user_id": self._sender_qq}
        else:
            return "小萌找不到要提醒谁 o_o"
        new_task = {
            "id": task_id,
            "type": task_type,
            "msg": f"⏰ 提醒：{message}",
            "target": target,
        }
        if task_type == "once":
            new_task["run_at"] = run_at
        else:
            new_task["cron"] = run_at
        tasks_file = self._data_dir / "skills" / "scheduler" / "tasks.json"
        tasks_file.parent.mkdir(parents=True, exist_ok=True)
        try:
            if tasks_file.exists():
                data = _json.loads(tasks_file.read_text(encoding="utf-8"))
            else:
                data = {"tasks": []}
            data.setdefault("tasks", []).append(new_task)
            tasks_file.write_text(_json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        except Exception as e:
            return f"保存提醒失败: {e}"
        if task_type == "once":
            human_time = run_at
        else:
            human_time = f"每天 {run_at}"
        return f"好的！小萌会在 {human_time} 提醒你：{message} ⏰"

    @staticmethod
    def _parse_reminder_time(text: str):
        """解析中文自然语言时间，返回 (task_type, run_at_str) 或 None。"""
        import re
        from datetime import datetime, timedelta
        now = datetime.now()
        text = text.strip()

        # X分钟后
        m = re.search(r"(\d+)\s*分钟后", text)
        if m:
            target = now + timedelta(minutes=int(m.group(1)))
            return ("once", target.strftime("%Y-%m-%d %H:%M"))

        # X小时后
        m = re.search(r"(\d+)\s*小时后", text)
        if m:
            target = now + timedelta(hours=int(m.group(1)))
            return ("once", target.strftime("%Y-%m-%d %H:%M"))

        # 每天/每日 + 时间
        if re.search(r"每[天日]", text):
            m = re.search(r"(\d+)\s*[点时](?:\s*(\d+)\s*分)?", text)
            if m:
                h = int(m.group(1))
                mi = int(m.group(2)) if m.group(2) else 0
                if re.search(r"下午|晚上|傍晚", text) and h < 12:
                    h += 12
                elif re.search(r"上午|早上|早", text) and h == 12:
                    h = 0
                return ("cron", f"{h:02d}:{mi:02d}")

        # 今天/今日/今晚/明天/后天 + 时间
        day_offset = 0
        if re.search(r"明天|明日", text):
            day_offset = 1
        elif re.search(r"后天|后日", text):
            day_offset = 2

        m = re.search(r"(\d+)\s*[点时](?:\s*(\d+)\s*分)?", text)
        if m:
            h = int(m.group(1))
            mi = int(m.group(2)) if m.group(2) else 0
            if re.search(r"下午|晚上|傍晚|夜", text) and h < 12:
                h += 12
            elif re.search(r"上午|早上|早晨|早", text) and h == 12:
                h = 0
            elif h < 7 and not re.search(r"上午|早", text):
                h += 12  # 凌晨1-6点不加，其他模糊小时数默认下午
            target = (now + timedelta(days=day_offset)).replace(hour=h, minute=mi, second=0, microsecond=0)
            if target <= now and day_offset == 0:
                target += timedelta(days=1)
            return ("once", target.strftime("%Y-%m-%d %H:%M"))

        return None


# ──────────────────────────────────────────────────────────────
# Skill 系统（重构后：基于 SkillRegistry）
# ──────────────────────────────────────────────────────────────

from .skills import (
    SkillRegistry, SkillExecutor, SkillContext,
    ModelTier, UserLevel,
)

_g_skill_registry_initialized = False


def init_skill_registry(skills_dir: str) -> SkillRegistry:
    """初始化技能注册中心（在 gateway 启动时调用一次）。"""
    global _g_skill_registry_initialized
    registry = SkillRegistry.get_instance()
    if not _g_skill_registry_initialized:
        registry.init_loader(skills_dir)
        _g_skill_registry_initialized = True
    return registry


def load_skills_prompt(skills_dir: str) -> str:
    """
    [向后兼容] 扫描 data/skills/ 下的技能文件，构建注入 system prompt 的技能段落。

    内部委托给 SkillRegistry + SkillExecutor，保留与旧代码的兼容性。
    """
    registry = init_skill_registry(skills_dir)
    skills = registry.list_enabled()
    if not skills:
        return ""

    index_lines = []
    body_sections = []
    for s in skills:
        emoji = s.emoji or ""
        label = f"{emoji} **{s.name}**" if emoji else f"- **{s.name}**"
        if s.description:
            label += f": {s.description}"
        index_lines.append(label)
        if s.body:
            body_sections.append(f"### {s.name}\n\n{s.body}")

    parts = ["## 可用技能\n"]
    if index_lines:
        parts.append("\n".join(index_lines))
    if body_sections:
        parts.append("\n\n---\n\n")
        parts.append("\n\n---\n\n".join(body_sections))

    return "\n".join(parts)


def build_context_skills_prompt(
    skills_dir: str,
    model_tier: ModelTier,
    user_level: UserLevel,
    identity: str = "",
) -> str:
    """上下文感知的技能 prompt：根据模型层级和用户权限过滤技能。"""
    registry = init_skill_registry(skills_dir)
    skills = registry.get_for_context(model_tier, user_level, identity)
    if not skills:
        return ""

    ctx = SkillContext(model_tier=model_tier, user_level=user_level, identity=identity)
    executor = SkillExecutor()
    return executor.build_active_prompt(skills, ctx)


def get_tool_schemas_for_context(
    base_tools: list,
    model_tier: ModelTier,
    user_level: UserLevel,
    identity: str = "",
    plugin_manager=None,
) -> list:
    """获取当前上下文可用的工具 schema 列表（内置工具 + 插件工具）。"""
    from .skills import ToolSchemaBuilder
    return ToolSchemaBuilder.from_tools(base_tools, model_tier, user_level, identity, plugin_manager)


# ──────────────────────────────────────────────────────────────
# 内部辅助
# ──────────────────────────────────────────────────────────────

def _ddg_search(query: str, max_results: int = 4, proxy: str = "") -> List[Dict]:
    """同步 DuckDuckGo 搜索，用 asyncio.to_thread 包装后调用。"""
    from ddgs import DDGS
    kwargs = {"proxy": proxy} if proxy else {}
    with DDGS(**kwargs) as ddgs:
        return list(ddgs.text(query, max_results=max_results))


def parse_tool_calls(response_content: str) -> List[Dict]:
    """
    兜底解析：如果模型没走 function-calling 而是在文本里输出了 JSON，
    尝试提取 tool_call 结构。
    正常情况下不会用到这个，OpenAI 格式的 tool_calls 由 adapter 直接解析。
    """
    results = []
    for m in re.finditer(r'\{"tool":\s*"(\w+)".*?\}', response_content, re.DOTALL):
        try:
            obj = json.loads(m.group())
            if "tool" in obj and "args" in obj:
                results.append({"name": obj["tool"], "arguments": obj["args"]})
        except Exception:
            pass
    return results
